import pandas as pd
import pdfplumber
from PIL import Image
import pytesseract
from pptx import Presentation
import io
from typing import Dict, List, Any
import numpy as np
from langchain_core.documents import Document
import os
import streamlit as st

def handle_pdf(file) -> tuple[List[Document], List[Dict]]:
    """Handle large PDF files"""
    documents = []
    tables = []
    chunk_size = 10  # Process 10 pages at a time
    
    with pdfplumber.open(file) as pdf:
        total_pages = len(pdf.pages)
        
        # Process PDF in chunks
        for chunk_start in range(0, total_pages, chunk_size):
            chunk_end = min(chunk_start + chunk_size, total_pages)
            
            # Show progress
            progress_text = f"Processing PDF pages {chunk_start + 1} to {chunk_end} of {total_pages}"
            with st.spinner(progress_text):
                for page_num in range(chunk_start, chunk_end):
                    page = pdf.pages[page_num]
                    
                    # Extract text with memory-efficient processing
                    try:
                        text = page.extract_text()
                        if text:
                            documents.append(Document(
                                page_content=text,
                                metadata={
                                    "source": file.name,
                                    "page": page_num + 1,
                                    "type": "pdf_text"
                                }
                            ))
                        
                        # Extract tables with memory management
                        for table in page.extract_tables():
                            if table:
                                # Clean and validate table data
                                cleaned_table = [
                                    [str(cell).strip() if cell else "" for cell in row]
                                    for row in table
                                ]
                                
                                tables.append({
                                    "raw_data": cleaned_table,
                                    "page": page_num + 1,
                                    "source": file.name,
                                    "type": "pdf_table"
                                })
                    except Exception as e:
                        st.warning(f"Error processing page {page_num + 1}: {str(e)}")
                        continue
    
    return documents, tables

def handle_excel(file) -> tuple[List[Document], List[Dict]]:
    """Handle large Excel files"""
    documents = []
    tables = []
    
    try:
        # Read Excel file
        excel_file = pd.ExcelFile(file)
        total_sheets = len(excel_file.sheet_names)
        
        for sheet_idx, sheet_name in enumerate(excel_file.sheet_names, 1):
            with st.spinner(f"Processing sheet {sheet_idx}/{total_sheets}: {sheet_name}"):
                try:
                    # Read the entire sheet
                    df = pd.read_excel(
                        file,
                        sheet_name=sheet_name
                    )
                    
                    if not df.empty:
                        # Convert DataFrame to string representation
                        text_content = f"Sheet: {sheet_name}\n"
                        text_content += df.to_string(index=False)
                        
                        # Add to documents
                        documents.append(Document(
                            page_content=text_content,
                            metadata={
                                "source": file.name,
                                "sheet": sheet_name,
                                "type": "excel_text",
                                "rows": len(df)
                            }
                        ))
                        
                        # Store table data
                        tables.append({
                            "raw_data": df.values.tolist(),
                            "headers": df.columns.tolist(),
                            "sheet": sheet_name,
                            "source": file.name,
                            "type": "excel_table",
                            "rows": len(df)
                        })
                        
                except Exception as e:
                    st.warning(f"Error processing sheet {sheet_name}: {str(e)}")
                    continue
                
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        return documents, tables
    
    return documents, tables

def handle_powerpoint(file) -> List[Document]:
    """Handle PowerPoint files"""
    documents = []
    
    presentation = Presentation(file)
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        if slide_text:
            documents.append(Document(
                page_content="\n".join(slide_text),
                metadata={
                    "source": file.name,
                    "slide": slide_num,
                    "type": "powerpoint"
                }
            ))
    
    return documents

def handle_image(file) -> tuple[List[Document], List[Dict]]:
    """Handle large images"""
    documents = []
    images = []
    
    try:
        # Open image
        with Image.open(file) as image:
            # Calculate new dimensions while maintaining aspect ratio
            max_dimension = 2000  # Maximum dimension for processing
            ratio = min(max_dimension / image.width, max_dimension / image.height)
            new_size = (int(image.width * ratio), int(image.height * ratio))
            
            # Resize image if necessary
            if ratio < 1:
                image = image.resize(new_size, Image.Resampling.LANCZOS)
            
            # Convert to RGB if necessary
            if image.mode not in ('L', 'RGB'):
                image = image.convert('RGB')
            
            # Process image in memory-efficient manner
            with st.spinner("Processing image with OCR..."):
                try:
                    # Set Tesseract path for Windows
                    if os.name == 'nt':
                        pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
                    
                    # Perform OCR with optimized settings
                    text = pytesseract.image_to_string(
                        image,
                        lang='eng',
                        config='--psm 3 --oem 3 --dpi 300'
                    ).strip()
                    
                    if text:
                        documents.append(Document(
                            page_content=text,
                            metadata={
                                "source": file.name,
                                "type": "image_text",
                                "dimensions": f"{image.width}x{image.height}"
                            }
                        ))
                    
                    # Save processed image
                    buffered = io.BytesIO()
                    image.save(buffered, format='JPEG', quality=85, optimize=True)
                    
                    images.append({
                        "image": image,
                        "text": text if text else "No text extracted",
                        "source": file.name,
                        "type": "image",
                        "dimensions": f"{image.width}x{image.height}"
                    })
                    
                except Exception as e:
                    st.warning(f"OCR processing error: {str(e)}")
                    images.append({
                        "image": image,
                        "text": "OCR processing failed",
                        "source": file.name,
                        "type": "image",
                        "dimensions": f"{image.width}x{image.height}"
                    })
    
    except Exception as e:
        raise RuntimeError(f"Error processing image: {str(e)}")
    
    return documents, images 